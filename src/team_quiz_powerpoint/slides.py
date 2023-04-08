import tempfile
from typing import Optional, Tuple, List

import qrcode
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Cm, Pt

from team_quiz_powerpoint.model import Question, Quiz, Player


def find_shape_by_id(shapes, shape_id):
    """Return shape by shape_id."""
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def rules_slide(prs: Presentation, rules: List[str], title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[7])
    shapes = slide.shapes

    shapes.title.text = title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    if len(rules) > 0:
        tf.text = rules[0]

    if len(rules) > 1:
        for i in range(1, len(rules)):
            p = tf.add_paragraph()
            p.text = rules[i]
            p.level = 0


def title_slide(prs: Presentation, title: str, music: bool = False):
    if music:
        slide = prs.slides.add_slide(prs.slide_layouts[8])
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    shapes.title.text = title


def qr_slide(prs: Presentation, url: str, title: str = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    if title:
        shapes.title.text = title

    img = qrcode.make(url)

    with tempfile.TemporaryFile() as fp:
        img.save(fp)
        shapes.add_picture(fp, Cm(8.2), Cm(3.8), width=Cm(9))


def question_slide(prs: Presentation, question: Question, index: int, timer: bool, music: bool):
    if timer:
        if music:
            layout = prs.slide_layouts[2]
        else:
            layout = prs.slide_layouts[9]
    else:
        layout = prs.slide_layouts[3]

    slide = prs.slides.add_slide(layout)
    find_shape_by_id(slide.shapes, 7).text = f"{index}. {question.question}"
    find_shape_by_id(slide.shapes, 2).text = f"A. {question.get_answers()[0]}"
    find_shape_by_id(slide.shapes, 3).text = f"B. {question.get_answers()[1]}"
    find_shape_by_id(slide.shapes, 4).text = f"C. {question.get_answers()[2]}"
    find_shape_by_id(slide.shapes, 5).text = f"D. {question.get_answers()[3]}"

    find_shape_by_id(slide.shapes, 6).text = f"{question.team.name.capitalize()}"


def podium_slide(prs: Presentation, player1: Optional[Tuple[Player, int]] = None,
                 player2: Optional[Tuple[str, int]] = None,
                 player3: Optional[Tuple[str, int]] = None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = f"Podium"

    if player1:
        find_shape_by_id(slide.shapes, 3).text = f"{player1[0].name} - {player1[1][0]}"
    else:
        find_shape_by_id(slide.shapes, 3).text = ""

    if player2:
        find_shape_by_id(slide.shapes, 4).text = f"{player2[0].name} - {player2[1][0]}"
    else:
        find_shape_by_id(slide.shapes, 4).text = ""

    if player3:
        find_shape_by_id(slide.shapes, 5).text = f"{player3[0].name} - {player3[1][0]}"
    else:
        find_shape_by_id(slide.shapes, 5).text = ""


def teams_slide(prs: Presentation, quiz: Quiz):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = f"Teams"

    chart_data = CategoryChartData()

    for question_team in quiz.teams:
        series = []
        for (answer_team, correct, incorrect) in quiz.get_team_scores(question_team):
            if (correct + incorrect) <= 0:
                continue
            chart_data.categories.add_category(answer_team.name.capitalize())
            series.append((correct / (correct + incorrect)))
        chart_data.add_series(f'{question_team.name.capitalize()} questions', series)

    x, y, cx, cy = Cm(1.98), Cm(3.46), Cm(21.44), Cm(9.51)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.value_axis.maximum_scale = 1
    chart.value_axis.tick_labels.number_format = '0%'
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False


def teams_winning_slide(prs: Presentation, quiz: Quiz):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = f"Teams"

    chart_data = CategoryChartData()

    series = []
    for (answer_team, correct, incorrect) in quiz.get_team_scores():
        if (correct + incorrect) <= 0:
            continue
        chart_data.categories.add_category(answer_team.name.capitalize())
        series.append((correct / (correct + incorrect)))
    chart_data.add_series(f'', series)

    if len(series) > 0:

        x, y, cx, cy = Cm(1.98), Cm(3.46), Cm(21.44), Cm(9.51)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.value_axis.maximum_scale = 1
        chart.value_axis.tick_labels.number_format = '0%'
        for series in chart.series:
            fill = series.format.fill  # fill the legend as well
            fill.solid()
            fill.fore_color.rgb = RGBColor(79, 239, 190)
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(13)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END


def answer_slide(prs: Presentation, quiz: Quiz, question_index: int, revealed: bool):
    question: Question = quiz.questions[question_index]

    layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(layout)
    find_shape_by_id(slide.shapes, 7).text = f"{question_index}. {question.question}"
    find_shape_by_id(slide.shapes, 2).text = f"A. {question.get_answers()[0]}"
    find_shape_by_id(slide.shapes, 3).text = f"B. {question.get_answers()[1]}"
    find_shape_by_id(slide.shapes, 4).text = f"C. {question.get_answers()[2]}"
    find_shape_by_id(slide.shapes, 5).text = f"D. {question.get_answers()[3]}"

    if revealed:
        correct_textbox = find_shape_by_id(slide.shapes, int(question.get_correct_answer()) + 1)
        correct_textbox.fill.solid()
        correct_textbox.fill.fore_color.rgb = RGBColor(79, 239, 190)

    if revealed:
        find_shape_by_id(slide.shapes, 6).text = f"{question.team.name.capitalize()} - {question.author}"
    else:
        find_shape_by_id(slide.shapes, 6).text = f"{question.team.name.capitalize()}"

    if revealed:
        scores = quiz.get_team_scores_for_question(question_index)

        chart_data = CategoryChartData()

        series = []
        for (team, correct, incorrect) in scores:
            if (correct + incorrect) <= 0:
                continue
            chart_data.categories.add_category(team.name.capitalize())
            series.append((correct / (correct + incorrect)))
        chart_data.add_series('', series)

        if len(series) > 0:
            x, y, cx, cy = Cm(19), Cm(3), Cm(5), Cm(10)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            chart.value_axis.maximum_scale = 1
            chart.value_axis.tick_labels.number_format = '0%'
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(79, 239, 190)
